import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Helper function to add a slide with a title and content
    def add_slide(title_text, subtitle_text=None, bullet_points=None):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title styling
        title = slide.shapes.title
        title.text = title_text
        
        # Content styling
        if bullet_points:
            tf = slide.placeholders[1].text_frame
            tf.text = bullet_points[0]
            for point in bullet_points[1:]:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
        elif subtitle_text:
            slide.placeholders[1].text = subtitle_text

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Campus Connect: IITP Super-App"
    subtitle.text = "Live Transit Tracking & Student Marketplace\nEngineering a Smarter Campus"

    # Slide 2: The Problem
    add_slide("The Challenge", bullet_points=[
        "Transit Blindness: Students have no real-time visibility into bus locations.",
        "Inefficient Commutes: Wasted time waiting at stops or missing buses.",
        "Fragmented Commerce: P2P selling is restricted to unorganized social media groups.",
        "Lack of Centralized Infrastructure: No single 'source of truth' for campus services."
    ])

    # Slide 3: The Solution
    add_slide("The Vision: Campus Connect", bullet_points=[
        "A Unified Super-App: Integrating logistics and commerce into one platform.",
        "Real-Time Awareness: Empowering students with live data.",
        "Secure P2P Ecosystem: Building trust within the campus community.",
        "Scalable Architecture: Designed to grow with IIT Patna's needs."
    ])

    # Slide 4: Feature - Live Transit Tracking
    add_slide("Live Transit Tracking", bullet_points=[
        "Live GPS Streams: Driver-side app sends high-frequency telemetry via WebSockets.",
        "Dead Reckoning Engine: Sophisticated interpolation logic estimates bus positions even during signal drops.",
        "Visual Intelligence: Map markers change state (Live/Estimated/Inactive) based on stream health.",
        "Integrated Schedules: Official IITP bus routes and timings baked in."
    ])

    # Slide 5: Feature - Student Marketplace
    add_slide("Student Marketplace", bullet_points=[
        "Peer-to-Peer Commerce: Direct listings for books, cycles, and electronics.",
        "Real-Time Socket Chat: Negotiation and communication happen instantly within the app.",
        "Secure Environment: Restricted to university members (IITP ecosystem).",
        "Rich UI: High-performance product discovery and management."
    ])

    # Slide 6: Tech Stack
    add_slide("The Tech Stack", bullet_points=[
        "Frontend: React 18, Vite, Tailwind CSS, Socket.io-client.",
        "Backend: Node.js, Express, Socket.io.",
        "Database & ORM: PostgreSQL with Prisma ORM.",
        "Infrastructure: Fully Dockerized (PostgreSQL + API).",
        "Mobile: Expo / React Native for Driver Telemetry."
    ])

    # Slide 7: System Architecture
    add_slide("Hybrid Architecture", bullet_points=[
        "Containerized Backend: Seamless deployment and consistency via Docker Compose.",
        "Bi-directional Sockets: Low-latency event-driven communication for live tracking and chat.",
        "Dead Reckoning Logic: Frontend-heavy interpolation to ensure smooth visual experience.",
        "Schema Safety: Strong typing and migrations via Prisma."
    ])

    # Slide 8: Future Roadmap
    add_slide("The Road Ahead", bullet_points=[
        "AI-Powered Demand Prediction: Estimating bus crowds based on historical data.",
        "Unified Payment Gateway: Integrating UPI for marketplace transactions.",
        "Club & Event Management: Expanding the super-app feature set.",
        "Campus-wide Notifications: Emergency and official broadcasts."
    ])

    # Save the presentation
    output_path = "Campus_Connect_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
